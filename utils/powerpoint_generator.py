"""
PowerPoint Generator Module
Creates PowerPoint presentations from AI-generated content
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, Any, List, Optional
import os


class PowerPointGenerator:
    """Generates PowerPoint presentations from structured content"""
    
    # Color schemes
    COLOR_SCHEMES = {
        "professional": {
            "primary": "1F4E79",      # Dark blue
            "secondary": "2E75B6",    # Medium blue
            "accent": "5B9BD5",       # Light blue
            "text": "333333",         # Dark gray
            "background": "FFFFFF"    # White
        },
        "academic": {
            "primary": "4A0E4E",      # Deep purple
            "secondary": "7B2D8E",    # Medium purple
            "accent": "9B59B6",       # Light purple
            "text": "2C3E50",         # Dark slate
            "background": "FFFFFF"
        },
        "creative": {
            "primary": "E74C3C",      # Red
            "secondary": "3498DB",    # Blue
            "accent": "F39C12",       # Orange
            "text": "2C3E50",         # Dark slate
            "background": "FFFFFF"
        },
        "minimal": {
            "primary": "2C3E50",      # Dark slate
            "secondary": "7F8C8D",    # Gray
            "accent": "1ABC9C",       # Teal
            "text": "333333",         # Dark gray
            "background": "FFFFFF"
        },
        "modern": {
            "primary": "00D4AA",      # Teal
            "secondary": "6C5CE7",    # Purple
            "accent": "FD79A8",       # Pink
            "text": "2D3436",         # Dark
            "background": "FFFFFF"
        }
    }
    
    def __init__(self):
        self.prs = None
        self.colors = self.COLOR_SCHEMES["professional"]
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def create_presentation(
        self, 
        content: Dict[str, Any],
        output_path: str
    ) -> Optional[str]:
        """
        Create a PowerPoint presentation from generated content
        
        Args:
            content: Dictionary with presentation structure
            output_path: Path to save the presentation
            
        Returns:
            Path to created file or None if failed
        """
        try:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
            self.prs.slide_height = Inches(7.5)
            
            # Set color scheme
            theme = content.get('theme_suggestion', 'professional').lower()
            if theme in self.COLOR_SCHEMES:
                self.colors = self.COLOR_SCHEMES[theme]
            
            # Custom colors from content if provided
            custom_colors = content.get('color_scheme', {})
            if custom_colors:
                if custom_colors.get('primary'):
                    self.colors['primary'] = custom_colors['primary'].lstrip('#')
                if custom_colors.get('secondary'):
                    self.colors['secondary'] = custom_colors['secondary'].lstrip('#')
                if custom_colors.get('accent'):
                    self.colors['accent'] = custom_colors['accent'].lstrip('#')
            
            # Create slides
            slides = content.get('slides', [])
            
            for slide_data in slides:
                slide_type = slide_data.get('slide_type', 'content').lower()
                
                if slide_type == 'title':
                    self._create_title_slide(slide_data, content)
                elif slide_type == 'summary' or slide_type == 'questions':
                    self._create_summary_slide(slide_data)
                else:
                    self._create_content_slide(slide_data)
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
            
            # Save presentation
            if not output_path.endswith('.pptx'):
                output_path += '.pptx'
            
            self.prs.save(output_path)
            return output_path
            
        except Exception as e:
            print(f"Error creating presentation: {e}")
            return None
    
    def _create_title_slide(self, slide_data: Dict[str, Any], content: Dict[str, Any]):
        """Create a title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background shape
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self._hex_to_rgb(self.colors['primary'])
        background.line.fill.background()
        
        # Add accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(5.5),
            self.prs.slide_width, Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors['accent'])
        accent_bar.line.fill.background()
        
        # Main title
        title = content.get('presentation_title', slide_data.get('title', 'Presentation'))
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle = content.get('presentation_subtitle', '')
        if not subtitle and slide_data.get('bullet_points'):
            subtitle = slide_data['bullet_points'][0]
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add speaker notes
        if slide_data.get('speaker_notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['speaker_notes']
    
    def _create_content_slide(self, slide_data: Dict[str, Any]):
        """Create a content slide with bullet points"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(self.colors['primary'])
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Content')
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Bullet points
        bullet_points = slide_data.get('bullet_points', [])
        if bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.6),
                Inches(11.833), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, point in enumerate(bullet_points):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                
                para.text = f"• {point}"
                para.font.size = Pt(24)
                para.font.color.rgb = self._hex_to_rgb(self.colors['text'])
                para.space_after = Pt(18)
                para.level = 0
        
        # Accent line at bottom
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.35),
            self.prs.slide_width, Inches(0.15)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self._hex_to_rgb(self.colors['accent'])
        accent_line.line.fill.background()
        
        # Slide number
        slide_num = slide_data.get('slide_number', '')
        if slide_num:
            num_box = slide.shapes.add_textbox(
                Inches(12.5), Inches(7),
                Inches(0.5), Inches(0.3)
            )
            num_frame = num_box.text_frame
            num_para = num_frame.paragraphs[0]
            num_para.text = str(slide_num)
            num_para.font.size = Pt(12)
            num_para.font.color.rgb = self._hex_to_rgb(self.colors['secondary'])
            num_para.alignment = PP_ALIGN.RIGHT
        
        # Add speaker notes
        if slide_data.get('speaker_notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['speaker_notes']
    
    def _create_summary_slide(self, slide_data: Dict[str, Any]):
        """Create a summary or questions slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Full background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self._hex_to_rgb(self.colors['secondary'])
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Summary')
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Content
        bullet_points = slide_data.get('bullet_points', [])
        if bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2),
                Inches(10.333), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, point in enumerate(bullet_points):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                
                # Use checkmark for summary, question mark for questions
                slide_type = slide_data.get('slide_type', '').lower()
                prefix = "✓ " if slide_type == 'summary' else "❓ " if slide_type == 'questions' else "• "
                
                para.text = f"{prefix}{point}"
                para.font.size = Pt(26)
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.space_after = Pt(20)
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.35),
            self.prs.slide_width, Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors['accent'])
        accent_bar.line.fill.background()
        
        # Add speaker notes
        if slide_data.get('speaker_notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['speaker_notes']


# Singleton instance
powerpoint_generator = PowerPointGenerator()
