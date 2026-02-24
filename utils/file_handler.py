"""
File Handler Module
Handles file downloads, parsing, and content extraction from various file types
"""
import os
import io
import aiohttp
from typing import Dict, Any, Optional, List
from pathlib import Path
import PyPDF2
import docx
import pandas as pd
from PIL import Image
import base64

from config import ALLOWED_FILE_EXTENSIONS, MAX_FILE_SIZE_MB


class FileHandler:
    """Handles file operations including download and content extraction"""
    
    def __init__(self, temp_dir: str = "temp_files"):
        """Initialize FileHandler with a temporary directory"""
        self.temp_dir = temp_dir
        os.makedirs(temp_dir, exist_ok=True)
    
    async def download_attachment(self, url: str, filename: str) -> Optional[str]:
        """
        Download a file from Discord attachment URL
        
        Args:
            url: The attachment URL
            filename: Original filename
            
        Returns:
            Path to downloaded file or None if failed
        """
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url) as response:
                    if response.status == 200:
                        filepath = os.path.join(self.temp_dir, filename)
                        content = await response.read()
                        
                        # Check file size
                        size_mb = len(content) / (1024 * 1024)
                        if size_mb > MAX_FILE_SIZE_MB:
                            return None
                        
                        with open(filepath, 'wb') as f:
                            f.write(content)
                        return filepath
            return None
        except Exception as e:
            print(f"Error downloading file: {e}")
            return None
    
    async def extract_content(self, filepath: str) -> Dict[str, Any]:
        """
        Extract content from a file based on its type
        
        Args:
            filepath: Path to the file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        filename = os.path.basename(filepath)
        extension = Path(filepath).suffix.lower()
        
        result = {
            "filename": filename,
            "extension": extension,
            "content": "",
            "is_binary": False,
            "error": None
        }
        
        try:
            if extension == '.pdf':
                result["content"] = self._extract_pdf(filepath)
            elif extension in ['.docx', '.doc']:
                result["content"] = self._extract_docx(filepath)
            elif extension in ['.xlsx', '.xls']:
                result["content"] = self._extract_excel(filepath)
            elif extension == '.csv':
                result["content"] = self._extract_csv(filepath)
            elif extension in ['.txt', '.md', '.py', '.js', '.json']:
                result["content"] = self._extract_text(filepath)
            elif extension in ['.png', '.jpg', '.jpeg', '.gif', '.webp']:
                result["content"] = await self._extract_image(filepath)
                result["is_binary"] = True
                result["image_data"] = self._get_image_base64(filepath)
                result["mime_type"] = self._get_mime_type(extension)
            elif extension in ['.pptx', '.ppt']:
                result["content"] = self._extract_pptx(filepath)
            else:
                result["error"] = f"Unsupported file type: {extension}"
                
        except Exception as e:
            result["error"] = str(e)
        
        return result
    
    def _extract_pdf(self, filepath: str) -> str:
        """Extract text content from PDF file"""
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return '\n\n'.join(text)
    
    def _extract_docx(self, filepath: str) -> str:
        """Extract text content from Word document"""
        doc = docx.Document(filepath)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                paragraphs.append(row_text)
        
        return '\n\n'.join(paragraphs)
    
    def _extract_excel(self, filepath: str) -> str:
        """Extract content from Excel file"""
        xl = pd.ExcelFile(filepath)
        content = []
        
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(xl, sheet_name=sheet_name)
            content.append(f"=== Sheet: {sheet_name} ===")
            content.append(df.to_string())
        
        return '\n\n'.join(content)
    
    def _extract_csv(self, filepath: str) -> str:
        """Extract content from CSV file"""
        df = pd.read_csv(filepath)
        return f"CSV Content:\n{df.to_string()}"
    
    def _extract_text(self, filepath: str) -> str:
        """Extract content from text file"""
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    async def _extract_image(self, filepath: str) -> str:
        """Get image info (actual analysis done by Gemini)"""
        with Image.open(filepath) as img:
            return f"[Image: {img.format}, Size: {img.size[0]}x{img.size[1]}, Mode: {img.mode}]"
    
    def _get_image_base64(self, filepath: str) -> bytes:
        """Get raw image bytes for Gemini Vision"""
        with open(filepath, 'rb') as f:
            return f.read()
    
    def _get_mime_type(self, extension: str) -> str:
        """Get MIME type for image extension"""
        mime_types = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.webp': 'image/webp'
        }
        return mime_types.get(extension, 'image/png')
    
    def _extract_pptx(self, filepath: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content.append(f"=== Slide {slide_num} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text)
            
            return '\n\n'.join(content)
        except ImportError:
            return "[PowerPoint content - python-pptx not installed]"
    
    def is_valid_file(self, filename: str) -> bool:
        """Check if file extension is allowed"""
        extension = Path(filename).suffix.lower()
        return extension in ALLOWED_FILE_EXTENSIONS
    
    def cleanup(self, filepath: Optional[str] = None):
        """Clean up temporary files"""
        if filepath and os.path.exists(filepath):
            os.remove(filepath)
        elif os.path.exists(self.temp_dir):
            for file in os.listdir(self.temp_dir):
                os.remove(os.path.join(self.temp_dir, file))


# Singleton instance
file_handler = FileHandler()
